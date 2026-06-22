"""Generate an office automation pack from the chapter CSV data."""
import csv
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path.cwd()
INPUT = ROOT / "input" / "report_data.csv"
REPORTS = ROOT / "reports"


def read_rows():
    with INPUT.open(encoding="utf-8") as f:
        return list(csv.DictReader(f))


def make_excel(rows):
    wb = Workbook()
    ws = wb.active
    ws.title = "card_factory"
    headers = ["chapter", "status", "artifact", "score"]
    ws.append(headers)
    for row in rows:
        ws.append([row["chapter"], row["status"], row["artifact"], int(row["score"])])

    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="2F6BFF")
    ws.column_dimensions["A"].width = 12
    ws.column_dimensions["B"].width = 12
    ws.column_dimensions["C"].width = 24
    ws.column_dimensions["D"].width = 10
    path = REPORTS / "final_report.xlsx"
    wb.save(path)
    return path


def make_word(rows):
    scores = [int(row["score"]) for row in rows]
    doc = Document()
    doc.add_heading("科研卡片工厂结课报告", level=1)
    doc.add_paragraph(f"本报告由 Python 自动生成，共汇总 {len(rows)} 个章节作品。")
    doc.add_paragraph(f"平均完成分：{sum(scores) / len(scores):.1f}")

    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    for i, header in enumerate(["章节", "状态", "作品", "完成分"]):
        table.rows[0].cells[i].text = header
    for row in rows:
        cells = table.add_row().cells
        cells[0].text = row["chapter"]
        cells[1].text = row["status"]
        cells[2].text = row["artifact"]
        cells[3].text = row["score"]

    doc.add_paragraph("复盘提示：自动化最重要的不是少打几个字，而是让报告结构稳定、结果可复现。")
    path = REPORTS / "final_report.docx"
    doc.save(path)
    return path


def make_ppt(rows):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "科研卡片工厂结课报告"
    title.text_frame.paragraphs[0].font.size = Pt(30)

    left = Inches(0.9)
    top = Inches(1.55)
    width = Inches(8.0)
    height = Inches(0.55)
    scores = [int(row["score"]) for row in rows]
    summary = f"汇总 {len(rows)} 个章节作品，平均完成分 {sum(scores) / len(scores):.1f}"
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text = summary
    box.text_frame.paragraphs[0].font.size = Pt(18)

    for i, row in enumerate(rows[:5]):
        item = slide.shapes.add_textbox(left, Inches(2.35 + i * 0.55), width, height)
        item.text = f"{row['chapter']} | {row['artifact']} | {row['score']}"
        item.text_frame.paragraphs[0].font.size = Pt(16)

    path = REPORTS / "final_slides.pptx"
    prs.save(path)
    return path


def font(size):
    candidates = [
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/arial.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def make_preview(rows):
    im = Image.new("RGB", (1500, 950), "#F7F8FB")
    d = ImageDraw.Draw(im)
    d.rounded_rectangle((90, 70, 1410, 880), radius=26, fill="#FFFFFF", outline="#D8E0EC", width=3)
    d.text((150, 125), "科研卡片工厂结课报告", fill="#162033", font=font(54))
    d.text((150, 205), "Python 自动汇总 CSV，生成 Markdown、Word、Excel 和 PPT。", fill="#5F6673", font=font(28))

    x0, y0 = 150, 300
    widths = [170, 170, 430, 140]
    headers = ["章节", "状态", "作品", "完成分"]
    colors = ["#2F6BFF", "#24A06B", "#F28C28", "#7A5AF8"]
    x = x0
    for header, width, color in zip(headers, widths, colors):
        d.rounded_rectangle((x, y0, x + width, y0 + 58), radius=12, fill=color)
        d.text((x + 22, y0 + 13), header, fill="#FFFFFF", font=font(26))
        x += width + 12

    for r, row in enumerate(rows):
        y = y0 + 88 + r * 76
        values = [row["chapter"], row["status"], row["artifact"], row["score"]]
        x = x0
        for value, width in zip(values, widths):
            d.rounded_rectangle((x, y, x + width, y + 54), radius=12, fill="#F1F5F9", outline="#E2E8F0", width=2)
            d.text((x + 20, y + 12), str(value), fill="#162033", font=font(24))
            x += width + 12

    path = REPORTS / "final_report_preview.png"
    im.save(path, optimize=True, quality=95)
    return path


def main():
    REPORTS.mkdir(exist_ok=True)
    rows = read_rows()
    generated = [
        make_excel(rows),
        make_word(rows),
        make_ppt(rows),
        make_preview(rows),
    ]
    print("已生成办公自动化成果包：")
    for path in generated:
        print(f"- {path}")


if __name__ == "__main__":
    main()
